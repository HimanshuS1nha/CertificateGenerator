from flask import Flask, render_template, request
import pandas as pd
import numpy as np
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import pythoncom
import win32com.client
import os
from flask import Flask
from flask_pymongo import PyMongo
import time
import qrcode
import random


def make_certificate(filename, template, date, title):
    i = 0
    num = int(time.time()*random.randint(1, 1000))
    data = pd.read_excel(filename)
    names = np.array(data['Name'])
    for name in names:
        prs = Presentation(template)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if (shape.text.find("{{FULL_NAME}}")) != -1:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                cur_text = run.text
                                i += 1
                                new_text = cur_text.replace(
                                    "{{FULL_NAME}}", str(name))
                                run.text = new_text
                    if (shape.text.find("{{DATE}}")) != -1:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                cur_text = run.text
                                i += 1
                                new_text = cur_text.replace(
                                    "{{DATE}}", str(date))
                                run.text = new_text
                    if (shape.text.find("{{QR_HERE}}")) != -1:
                        qr = qrcode.QRCode(version=None, box_size=4)
                        data = f"http://localhost:3000/api/verify/{num+i}"
                        qr.add_data(data)
                        qr.make(fit=True)
                        img = qr.make_image(fill_color='black',
                                            back_color='white')
                        img.save(f"static/{name}qrcode.png")
                        img_path = f"static/{name}qrcode.png"
                        slide.shapes.add_picture(
                            img_path, 0, 0)
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "{{QR_HERE}}" in run.text:
                                    # replace "NAME" with an empty string
                                    run.text = run.text.replace(
                                        "{{QR_HERE}}", "")

                # if shape.has_text_frame:
                #     if (shape.text.find("{{SIGN}}")) != -1:
                #         img_path = sign
                #         # left = Inches(shape.left)
                #         # top = Inches(shape.top)
                #         slide.shapes.add_picture(
                #             img_path, 0, 0)
                #         text_frame = shape.text_frame
                #         for paragraph in text_frame.paragraphs:
                #             for run in paragraph.runs:
                #                 if "{{SIGN}}" in run.text:
                #                     run.text = run.text.replace(
                #                         "{{SIGN}}", "")

                # qr = qrcode.QRCode(version=None, box_size=1)
                # data = f"https://localhost:3000/api/verify/{num+i}"
                # qr.add_data(data)
                # qr.make(fit=True)
                # img = qr.make_image(fill_color='black',
                #                     back_color='white')
                # img.save(f"static/{name}qrcode.png")
                # img_path = f"static/{name}qrcode.png"
                # slide.shapes.add_picture(
                #     img_path, 0, 0)

                    # pic = slide.shapes.add_picture(
                    #     img_path, 0, 0)
                    prs.save(f"static/{name}.pptx")
                    i += 1
        certificate = {'name': name,
                       'Occasion': title,
                       'date': date,
                       'cid': num+i}
        mongo.db.certificates.insert_one(certificate)


def convert():
    Path = os.listdir(os.getcwd()+'\\static')
    folder = 'static'
    for file in Path:
        if file.endswith('.pptx'):
            output_file = os.path.splitext(file)[0] + '.pdf'

            # Use pywin32 to convert the PowerPoint file to a PDF
            pythoncom.CoInitialize()
            ppt = win32com.client.Dispatch('PowerPoint.Application')
            presentation = ppt.Presentations.Open(
                os.path.abspath(os.path.join(folder, file)))
            presentation.SaveAs(os.path.abspath(
                os.path.join(folder, output_file)), 32)
            presentation.Close()
            ppt.Quit()
            pythoncom.CoUninitialize()
            os.remove(os.path.join(folder, file))


def send_mail(filename, title):
    data = pd.read_csv(filename)
    names = np.array(data['Name'])
    emails = np.array(data['Email'])

    for i in range(len(emails)):
        fromaddr = "EMAILID"
        toaddr = emails[i]

        msg = MIMEMultipart()

        msg['From'] = fromaddr

        msg['To'] = toaddr

        msg['Subject'] = f"Certificate for {title}"

        body = f"Hello {names[i]}, Here is your certificate for {title}"

        msg.attach(MIMEText(body, 'plain'))

        filename = f"static/{names[i]}.pdf"
        attachment = open(filename, "rb")

        p = MIMEBase('application', 'octet-stream')

        p.set_payload((attachment).read())

        encoders.encode_base64(p)

        p.add_header('Content-Disposition',
                     "attachment; filename= %s" % filename)

        msg.attach(p)

        s = smtplib.SMTP('smtp.gmail.com', 587)

        s.starttls()

        s.login('EMAILID', "APP PASSWORD")

        text = msg.as_string()

        s.sendmail(fromaddr, toaddr, text)

        s.quit()


app = Flask(__name__)
# MongoDB URI
app.config["MONGO_URI"] = "MONGODB URI"
mongo = PyMongo(app)  # Initialising MongoDB


@app.route("/")
def home():
    return render_template("index.html")


@app.route('/api/<cid>')
def hello(cid):
    certificate = mongo.db.certificates.find_one({'cid': cid})
    if (certificate):
        # Certificate is verified
        return render_template('verify.html')
    else:
        # Certificate is not verified
        return render_template("not-verify.html")


@app.route("/api/file", methods=["GET", "POST"])
def operation():
    if request.method == "POST":
        date = request.form.get('date')  # Get the date from the form
        title = request.form.get('title')  # Get the title from the form
        f1 = request.files['file']  # Get the csv file from the form
        # Get the certificate template from the form
        f2 = request.files['template']
        # f3 = request.files['sign']
        f1.save(f'uploads/{secure_filename(f1.filename)}')
        f2.save(f'uploads/{secure_filename(f2.filename)}')
        # f3.save(f'uploads/{secure_filename(f3.filename)}')
        make_certificate(f1.filename, f2.filename, date, title)
        convert()  # For converting the output pptx file to pdf file
        # send_mail(f1.filename, title)  # For sending mail
        return render_template("results.html")


app.run(debug=True)
